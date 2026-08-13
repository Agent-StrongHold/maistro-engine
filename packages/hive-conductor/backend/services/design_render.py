"""DesignRenderService — Server-side rendering backends for design outputs.

Implements actual rendering for PDF (weasyprint), PPTX (python-pptx), DOCX (python-docx).
Uses background job queue (placeholder for Phase 2B.2).

Phase 2B.1: Render implementations
Phase 2B.2: Background job queue + S3 storage
"""

from __future__ import annotations

import html
import json
import logging
from io import BytesIO
from typing import Any

logger = logging.getLogger("hive.design_render")

__all__ = ["DesignRenderService"]


class DesignRenderService:
    """Server-side rendering implementations for design outputs."""

    def __init__(self) -> None:
        """Initialize render service."""
        self._render_cache: dict[str, bytes] = {}

    def _parse_markdown_lines(self, doc: Any, content: str) -> None:
        """Parse markdown-like content and add to document.

        Supports basic markdown: headings (# ## ###), bullet lists (-), paragraphs.
        """
        for line in content.split("\n"):
            if line.startswith("# "):
                doc.add_heading(line[2:], 1)
            elif line.startswith("## "):
                doc.add_heading(line[3:], 2)
            elif line.startswith("### "):
                doc.add_heading(line[4:], 3)
            elif line.startswith("- "):
                doc.add_paragraph(line[2:], style="List Bullet")
            elif line.strip():
                doc.add_paragraph(line)

    async def render_to_pdf(self, content: str, metadata: dict[str, Any]) -> bytes:
        """Render HTML/React to PDF via weasyprint.

        Args:
            content: HTML or rendered React output
            metadata: Output metadata (title, author, etc.)

        Returns:
            PDF bytes
        """
        try:
            from weasyprint import HTML

            # Wrap content in minimal HTML structure if not already HTML
            if not content.strip().startswith("<"):
                escaped_content = html.escape(content)
                escaped_title = html.escape(metadata.get("title", "Document"))
                html_content = f"""
                <!DOCTYPE html>
                <html>
                <head>
                    <meta charset="utf-8">
                    <title>{escaped_title}</title>
                    <style>
                        body {{
                            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
                            margin: 2cm;
                        }}
                        h1, h2, h3 {{ color: #333; }}
                    </style>
                </head>
                <body>
                    {escaped_content}
                </body>
                </html>
                """
            else:
                html_content = content

            # Render to PDF
            pdf_bytes = HTML(string=html_content).write_pdf()
            logger.info("PDF render successful (%d bytes)", len(pdf_bytes))
            return pdf_bytes

        except ImportError:
            logger.error("weasyprint not installed - PDF rendering unavailable")
            raise RuntimeError("PDF rendering requires weasyprint package") from None
        except Exception as exc:
            logger.error("PDF render failed: %s", exc)
            raise

    async def render_to_pptx(self, content: str, metadata: dict[str, Any]) -> bytes:
        """Render to PPTX via python-pptx.

        Expects content as JSON slides array or markdown.

        Args:
            content: Slide content (JSON or markdown)
            metadata: Presentation metadata (title, author, etc.)

        Returns:
            PPTX bytes
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            # Parse slides from content
            slides_data = []
            try:
                slides_data = json.loads(content)
            except json.JSONDecodeError:
                # Fallback: treat content as single slide
                slides_data = [{"title": metadata.get("title", "Slide"), "content": content}]

            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = metadata.get("title", "Presentation")
            subtitle.text = metadata.get("subtitle", "")

            # Add content slides
            for slide_data in slides_data:
                if isinstance(slide_data, dict):
                    blank_layout = prs.slide_layouts[6]  # Blank layout
                    slide = prs.slides.add_slide(blank_layout)

                    # Add title
                    if "title" in slide_data:
                        title_box = slide.shapes.add_textbox(
                            Inches(0.5), Inches(0.5), Inches(9), Inches(1)
                        )
                        title_frame = title_box.text_frame
                        title_frame.text = slide_data["title"]
                        title_frame.paragraphs[0].font.size = Pt(44)
                        title_frame.paragraphs[0].font.bold = True

                    # Add content
                    if "content" in slide_data:
                        content_box = slide.shapes.add_textbox(
                            Inches(0.5), Inches(1.75), Inches(9), Inches(5.25)
                        )
                        text_frame = content_box.text_frame
                        text_frame.word_wrap = True
                        text_frame.text = str(slide_data["content"])

            # Write to bytes
            output = BytesIO()
            prs.save(output)
            pptx_bytes = output.getvalue()
            logger.info("PPTX render successful (%d bytes)", len(pptx_bytes))
            return pptx_bytes

        except ImportError:
            logger.error("python-pptx not installed - PPTX rendering unavailable")
            raise RuntimeError("PPTX rendering requires python-pptx package") from None
        except Exception as exc:
            logger.error("PPTX render failed: %s", exc)
            raise

    async def render_to_docx(self, content: str, metadata: dict[str, Any]) -> bytes:
        """Render to DOCX via python-docx.

        Expects content as markdown or plain text.

        Args:
            content: Document content (markdown or plain text)
            metadata: Document metadata (title, author, etc.)

        Returns:
            DOCX bytes
        """
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = Document()

            # Add title
            if "title" in metadata:
                title = doc.add_heading(metadata["title"], 0)
                title.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Add metadata
            if "author" in metadata:
                doc.add_paragraph(f"Author: {metadata['author']}")
            if "created_at" in metadata:
                doc.add_paragraph(f"Created: {metadata['created_at']}")

            doc.add_paragraph()  # Spacing

            # Add content via markdown-like parsing (Phase 2B.3: full markdown support)
            self._parse_markdown_lines(doc, content)

            # Write to bytes
            output = BytesIO()
            doc.save(output)
            docx_bytes = output.getvalue()
            logger.info("DOCX render successful (%d bytes)", len(docx_bytes))
            return docx_bytes

        except ImportError:
            logger.error("python-docx not installed - DOCX rendering unavailable")
            raise RuntimeError("DOCX rendering requires python-docx package") from None
        except Exception as exc:
            logger.error("DOCX render failed: %s", exc)
            raise

    async def render_to_png(self, content: str, metadata: dict[str, Any]) -> bytes:
        """Render HTML to PNG via Playwright — not implemented yet.

        Phase 2B.3 will implement this. Until then it raises an
        `HTTPException(501)` rather than a bare `NotImplementedError`, so the
        caller gets an honest "this backend does not implement PNG" instead of
        a 500 that reads like a crash. Raised at the service boundary because
        `routes/design.py` re-raises `HTTPException` untouched and funnels every
        other exception into a generic 500 — a domain error would need each
        current and future PNG-dispatching route to translate it identically.

        Args:
            content: HTML content
            metadata: Image metadata (width, height, etc.)

        Raises:
            HTTPException: always, with status 501.
        """
        from fastapi import HTTPException

        logger.warning("PNG rendering not yet implemented (Phase 2B.3)")
        raise HTTPException(
            status_code=501,
            detail=(
                "PNG rendering is not implemented. It is deferred to Phase 2B.3 "
                "(Playwright renderer). Use pdf, pptx or docx instead."
            ),
        )


# Singleton instance
_singleton: DesignRenderService | None = None


def get_design_render_service() -> DesignRenderService:
    """Get the DesignRenderService singleton."""
    if _singleton is None:
        raise RuntimeError("DesignRenderService not initialized")
    return _singleton


def init_design_render_service() -> DesignRenderService:
    """Initialize the DesignRenderService singleton."""
    global _singleton
    if _singleton is None:
        _singleton = DesignRenderService()
    return _singleton
