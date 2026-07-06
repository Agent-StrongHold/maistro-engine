"""Canvas structured exporters — SPEC-070426-457b / ADR-070426-f2a0.

Serialize a fixed-layout page — an ordered stack of positioned image/text layers —
into **editable PPTX** and **absolute-positioned single-file HTML**. Text stays text
(a real PPTX text box / an HTML text node), never rasterized, which is what makes the
PPTX output editable in PowerPoint. This is the plugin-free "fixed-page" floor: pure
serialization over the layer tree, no LLM and no image backend.

HTML export needs only the standard library. PPTX export lazily imports ``python-pptx``
and raises :class:`ExporterDependencyError` when it is not installed, so the deck path
degrades gracefully rather than crashing an install that never uses it.

The export model deliberately mirrors ``CanvasRecord`` / ``LayerRecord`` / ``TextConfig``
field names (``x``, ``y``, ``z_index``, ``rotation``, ``opacity``; text ``font`` / ``size``
/ ``color`` / ``weight`` / ``alignment``) so adapting real canvas records is a shallow map.
"""

from __future__ import annotations

import base64
import html
from collections.abc import Sequence
from dataclasses import dataclass, field
from io import BytesIO
from typing import Any, cast

from maistro_canvas.types import CanvasError

_PX_TO_EMU = 9525  # 1 CSS px at 96 dpi
_PX_TO_PT = 0.75  # 1 CSS px at 96 dpi -> points


class ExporterDependencyError(CanvasError):
    """An optional export backend (e.g. python-pptx) is not installed."""

    code = "EXPORTER_DEPENDENCY_MISSING"


@dataclass
class ExportText:
    """A text run on a layer. Mirrors ``maistro_canvas.types.TextConfig``."""

    content: str
    font: str = "sans-serif"
    size: int = 48  # px
    color: str = "#111111"  # hex
    weight: str = "normal"  # normal | bold
    alignment: str = "left"  # left | center | right


@dataclass
class ExportLayer:
    """One positioned layer. Exactly one of ``image_png`` / ``image_src`` / ``text`` is set."""

    z_index: int = 0
    x: float = 0.0
    y: float = 0.0
    width: float | None = None
    height: float | None = None
    rotation: float = 0.0
    opacity: float = 1.0
    image_png: bytes | None = None  # binary — required for PPTX pictures
    image_src: str | None = None  # URL / data-uri — HTML only
    text: ExportText | None = None


@dataclass
class ExportPage:
    """A single fixed-size page. A deck is a sequence of these."""

    width: int
    height: int
    background_color: str = "#FFFFFF"
    layers: list[ExportLayer] = field(default_factory=list)


def _ordered(layers: Sequence[ExportLayer]) -> list[ExportLayer]:
    """Back-to-front paint order."""
    return sorted(layers, key=lambda layer: layer.z_index)


def _data_uri(png: bytes) -> str:
    return "data:image/png;base64," + base64.b64encode(png).decode("ascii")


# ─── HTML ────────────────────────────────────────────────────────────────────


def _layer_html(layer: ExportLayer) -> str:
    css = ["position:absolute", f"left:{layer.x:g}px", f"top:{layer.y:g}px"]
    if layer.width is not None:
        css.append(f"width:{layer.width:g}px")
    if layer.height is not None:
        css.append(f"height:{layer.height:g}px")
    if layer.rotation:
        css.append(f"transform:rotate({layer.rotation:g}deg)")
    if layer.opacity != 1.0:
        css.append(f"opacity:{layer.opacity:g}")

    src = layer.image_src or (_data_uri(layer.image_png) if layer.image_png else None)
    if src is not None:
        style = ";".join([*css, "object-fit:contain"])
        return f'<img style="{style}" src="{html.escape(src, quote=True)}" alt="" />'
    if layer.text is not None:
        t = layer.text
        css += [
            f"font-family:{html.escape(t.font, quote=True)}",
            f"font-size:{t.size:g}px",
            f"color:{html.escape(t.color, quote=True)}",
            f"font-weight:{html.escape(t.weight, quote=True)}",
            f"text-align:{html.escape(t.alignment, quote=True)}",
        ]
        return f'<div style="{";".join(css)}">{html.escape(t.content)}</div>'
    return ""  # empty layer — nothing to paint


def export_html(page: ExportPage) -> str:
    """Render a page to a single self-contained HTML string (absolute-positioned layers)."""
    body = "\n    ".join(_layer_html(layer) for layer in _ordered(page.layers))
    container_style = (
        f"position:relative;width:{page.width}px;height:{page.height}px;"
        f"background:{html.escape(page.background_color, quote=True)};overflow:hidden"
    )
    return (
        "<!doctype html>\n"
        '<html><head><meta charset="utf-8" />'
        '<meta name="viewport" content="width=device-width, initial-scale=1" />'
        "<style>*{margin:0;box-sizing:border-box}</style></head>\n"
        f'<body>\n  <div style="{container_style}">\n    {body}\n  </div>\n</body></html>\n'
    )


# ─── PPTX ────────────────────────────────────────────────────────────────────


def export_pptx(pages: Sequence[ExportPage]) -> bytes:
    """Render pages to an editable ``.pptx`` (one slide per page). Text layers become
    real, editable text boxes; image layers become pictures. Raises
    :class:`ExporterDependencyError` if ``python-pptx`` is not installed."""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Emu, Pt
    except ImportError as exc:  # graceful: the deck backend simply isn't available
        msg = "PPTX export requires python-pptx (pip install 'maistro-canvas[export]')"
        raise ExporterDependencyError(msg) from exc

    if not pages:
        msg = "export_pptx requires at least one page"
        raise ExporterDependencyError(msg)

    # python-pptx ships partial types; route its untyped constructors through Any so this
    # type-checks identically whether or not the package is installed in the mypy env.
    rgb = cast("Any", RGBColor)
    align = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

    prs = Presentation()
    prs.slide_width = Emu(pages[0].width * _PX_TO_EMU)
    prs.slide_height = Emu(pages[0].height * _PX_TO_EMU)
    blank = prs.slide_layouts[6]

    for page in pages:
        slide = prs.slides.add_slide(blank)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb.from_string(page.background_color.lstrip("#"))

        for layer in _ordered(page.layers):
            left, top = Emu(int(layer.x * _PX_TO_EMU)), Emu(int(layer.y * _PX_TO_EMU))
            width = Emu(int(layer.width * _PX_TO_EMU)) if layer.width else None
            height = Emu(int(layer.height * _PX_TO_EMU)) if layer.height else None

            if layer.image_png is not None:
                slide.shapes.add_picture(BytesIO(layer.image_png), left, top, width, height)
            elif layer.text is not None:
                t = layer.text
                box = slide.shapes.add_textbox(
                    left,
                    top,
                    width or Emu(page.width * _PX_TO_EMU),
                    height or Emu(t.size * 2 * _PX_TO_EMU),
                )
                para = box.text_frame.paragraphs[0]
                para.alignment = align.get(t.alignment, PP_ALIGN.LEFT)
                run = para.add_run()
                run.text = t.content
                run.font.size = Pt(t.size * _PX_TO_PT)
                run.font.bold = t.weight == "bold"
                run.font.name = t.font
                run.font.color.rgb = rgb.from_string(t.color.lstrip("#"))

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
